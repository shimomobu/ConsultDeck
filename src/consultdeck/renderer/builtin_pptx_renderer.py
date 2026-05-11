from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from consultdeck.models.slide_spec import LayoutType, Slide, SlideSpec
from consultdeck.models.template_spec import TemplateSpec


class BuiltinPptxRenderer:
    def render(
        self,
        spec: SlideSpec,
        template: TemplateSpec,
        output_dir: Path,
    ) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{spec.deck_id}.pptx"

        presentation = Presentation()
        for slide in spec.slides:
            self._add_slide(presentation, slide)

        presentation.save(output_path)
        return output_path

    def _add_slide(self, presentation: Presentation, slide_spec: Slide) -> None:
        if slide_spec.layout_type is LayoutType.TITLE:
            slide = self._add_title_slide(presentation, slide_spec)
        elif slide_spec.layout_type is LayoutType.CONTENT:
            slide = self._add_content_slide(presentation, slide_spec)
        elif slide_spec.layout_type is LayoutType.TWO_COLUMN:
            slide = self._add_two_column_slide(presentation, slide_spec)
        else:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._write_notes(slide, slide_spec.notes)

    def _add_title_slide(self, presentation: Presentation, slide_spec: Slide):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
        frame = box.text_frame
        frame.text = slide_spec.title
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.runs[0].font.size = Pt(36)

        if slide_spec.message:
            subtitle = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(3.8),
                Inches(7),
                Inches(0.8),
            )
            subtitle.text_frame.text = slide_spec.message
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return slide

    def _add_content_slide(self, presentation: Presentation, slide_spec: Slide):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_title(slide, slide_spec.title)
        self._add_bullets(slide, slide_spec.bullets, Inches(1), Inches(1.7), Inches(8))
        return slide

    def _add_two_column_slide(
        self,
        presentation: Presentation,
        slide_spec: Slide,
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._add_title(slide, slide_spec.title)

        left_items = slide_spec.bullets[::2]
        right_items = slide_spec.bullets[1::2]
        if not right_items and left_items:
            right_items = [slide_spec.message]

        self._add_bullets(slide, left_items, Inches(0.8), Inches(1.7), Inches(4.0))
        self._add_bullets(slide, right_items, Inches(5.0), Inches(1.7), Inches(4.0))
        return slide

    def _add_title(self, slide, title: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.7))
        frame = box.text_frame
        frame.text = title
        frame.paragraphs[0].runs[0].font.size = Pt(28)

    def _add_bullets(self, slide, bullets: list[str], left, top, width) -> None:
        box = slide.shapes.add_textbox(left, top, width, Inches(4.8))
        frame = box.text_frame
        frame.clear()

        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(18)

    def _write_notes(self, slide, notes: str | None) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
