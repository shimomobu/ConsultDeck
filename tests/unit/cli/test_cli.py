import os
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from consultdeck.__main__ import main
from consultdeck.llm.provider_factory import build_llm_provider


PROJECT_ROOT = Path(__file__).resolve().parents[3]


def _console_script() -> str:
    return str(Path(sys.executable).with_name("consultdeck"))


def _env() -> dict[str, str]:
    env = os.environ.copy()
    env["PYTHONPATH"] = str(PROJECT_ROOT / "src")
    return env


def test_console_script_help_uses_console_command_name() -> None:
    result = subprocess.run(
        [_console_script(), "--help"],
        cwd=PROJECT_ROOT,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode == 0
    assert result.stdout.startswith("usage: consultdeck ")
    assert "usage: python -m consultdeck" not in result.stdout


def test_module_help_uses_module_command_name() -> None:
    result = subprocess.run(
        [sys.executable, "-m", "consultdeck", "--help"],
        cwd=PROJECT_ROOT,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode == 0
    assert result.stdout.startswith("usage: python -m consultdeck ")


def test_main_help_path_is_covered(capsys) -> None:
    with pytest.raises(SystemExit) as exc_info:
        main(["--help"])

    assert exc_info.value.code == 0
    assert capsys.readouterr().out.startswith("usage: ")


def test_cli_generates_pptx(tmp_path) -> None:
    output_dir = tmp_path / "output"

    result = subprocess.run(
        [
            sys.executable,
            "-m",
            "consultdeck",
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "経営層",
            "--slides",
            "5",
            "--output",
            str(output_dir),
            "--deck-id",
            "deck-cli-test",
        ],
        cwd=PROJECT_ROOT,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode == 0
    output_path = output_dir / "deck-cli-test.pptx"
    assert str(output_path) in result.stdout
    assert output_path.exists()
    assert len(Presentation(output_path).slides) == 5


def test_main_generates_pptx(tmp_path, capsys) -> None:
    output_dir = tmp_path / "output"

    exit_code = main(
        [
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "経営層",
            "--slides",
            "3",
            "--output",
            str(output_dir),
            "--deck-id",
            "deck-main-test",
        ]
    )

    output_path = output_dir / "deck-main-test.pptx"
    assert exit_code == 0
    assert capsys.readouterr().out.strip() == str(output_path)
    assert output_path.exists()


def test_main_accepts_fake_llm_provider_option(tmp_path, capsys) -> None:
    output_dir = tmp_path / "output"

    exit_code = main(
        [
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "経営層",
            "--slides",
            "3",
            "--output",
            str(output_dir),
            "--deck-id",
            "deck-fake-llm",
            "--llm-provider",
            "fake",
        ]
    )

    output_path = output_dir / "deck-fake-llm.pptx"
    assert exit_code == 0
    assert capsys.readouterr().out.strip() == str(output_path)
    assert output_path.exists()
    presentation = Presentation(output_path)
    slide_text = "\n".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    assert "Fake provider bullet" in slide_text
    assert "課題の要点を確認する" not in slide_text


def test_llm_provider_factory_returns_none_by_default() -> None:
    assert build_llm_provider("none") is None


def test_llm_provider_factory_builds_fake_provider() -> None:
    assert build_llm_provider("fake") is not None


def test_cli_uses_default_templates_outside_project_root(tmp_path) -> None:
    cwd = tmp_path / "outside"
    cwd.mkdir()
    output_dir = tmp_path / "output"

    result = subprocess.run(
        [
            _console_script(),
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "経営層",
            "--slides",
            "3",
            "--output",
            str(output_dir),
            "--deck-id",
            "deck-default-template",
        ],
        cwd=cwd,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode == 0
    assert (output_dir / "deck-default-template.pptx").exists()


def test_cli_prefers_explicit_templates_over_default(tmp_path) -> None:
    cwd = tmp_path / "outside"
    cwd.mkdir()
    output_dir = tmp_path / "output"
    template_dir = tmp_path / "templates"
    template_dir.mkdir()
    (template_dir / "custom.yaml").write_text(
        """
template_id: custom_proposal
name: Custom Proposal
doc_type: proposal
use_case: 提案書
audience: 現場
phase: proposal
slide_structure:
  - 表紙
  - 課題
layout_rules:
  default: content
style_rules:
  font: Arial
output_targets:
  - pptx
""".lstrip(),
        encoding="utf-8",
    )

    result = subprocess.run(
        [
            _console_script(),
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "現場",
            "--slides",
            "2",
            "--output",
            str(output_dir),
            "--templates",
            str(template_dir),
            "--deck-id",
            "deck-explicit-template",
        ],
        cwd=cwd,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    output_path = output_dir / "deck-explicit-template.pptx"
    assert result.returncode == 0
    assert output_path.exists()
    assert len(Presentation(output_path).slides) == 2


def test_cli_fails_when_required_arguments_are_missing(tmp_path) -> None:
    result = subprocess.run(
        [sys.executable, "-m", "consultdeck", "--output", str(tmp_path)],
        cwd=PROJECT_ROOT,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode != 0
    assert "error:" in result.stderr


def test_cli_fails_when_template_does_not_match(tmp_path) -> None:
    result = subprocess.run(
        [
            sys.executable,
            "-m",
            "consultdeck",
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "現場",
            "--slides",
            "5",
            "--output",
            str(tmp_path),
        ],
        cwd=PROJECT_ROOT,
        env=_env(),
        text=True,
        capture_output=True,
        check=False,
    )

    assert result.returncode != 0
    assert "No matching template" in result.stderr


def test_main_fails_when_template_does_not_match(tmp_path, capsys) -> None:
    exit_code = main(
        [
            "--topic",
            "生成AIの業務活用",
            "--purpose",
            "proposal",
            "--audience",
            "現場",
            "--slides",
            "3",
            "--output",
            str(tmp_path),
        ]
    )

    captured = capsys.readouterr()
    assert exit_code == 1
    assert captured.out == ""
    assert "No matching template" in captured.err
