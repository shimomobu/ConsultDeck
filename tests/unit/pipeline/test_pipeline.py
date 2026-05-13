from pathlib import Path
from typing import get_type_hints

import pytest
from pptx import Presentation

from consultdeck.models.requirement_spec import RequirementSpec
from consultdeck.models.slide_spec import SlideSpec
from consultdeck.pipeline.pipeline import Pipeline, PipelineError
from consultdeck.renderer.base import Renderer
from consultdeck.slide.content_generator import (
    FakeLlmProvider,
    GeneratedSlideContent,
    LlmGenerationResult,
)
from consultdeck.slide.builder import SlideBuilder


PROJECT_ROOT = Path(__file__).resolve().parents[3]
TEMPLATE_DIR = PROJECT_ROOT / "assets" / "templates"


def _requirement(
    purpose: str = "proposal",
    audience: str = "経営層",
    slide_count: int = 5,
) -> RequirementSpec:
    return RequirementSpec(
        theme="生成AIの業務活用",
        purpose=purpose,
        audience=audience,
        slide_count=slide_count,
    )


def test_pipeline_generates_pptx(tmp_path) -> None:
    pipeline = Pipeline(template_dir=TEMPLATE_DIR)

    output_path = pipeline.run(_requirement(), tmp_path)

    assert output_path.exists()
    assert output_path.suffix == ".pptx"


def test_pipeline_creates_output_dir(tmp_path) -> None:
    pipeline = Pipeline(template_dir=TEMPLATE_DIR)
    output_dir = tmp_path / "nested" / "output"

    output_path = pipeline.run(_requirement(), output_dir)

    assert output_dir.exists()
    assert output_path.parent == output_dir


def test_pipeline_uses_requested_slide_count(tmp_path) -> None:
    pipeline = Pipeline(template_dir=TEMPLATE_DIR)

    output_path = pipeline.run(_requirement(slide_count=4), tmp_path)

    presentation = Presentation(output_path)
    assert len(presentation.slides) == 4


def test_pipeline_fails_when_no_template_matches(tmp_path) -> None:
    pipeline = Pipeline(template_dir=TEMPLATE_DIR)

    with pytest.raises(PipelineError, match="No matching template"):
        pipeline.run(_requirement(purpose="proposal", audience="現場"), tmp_path)


def test_pipeline_accepts_fixed_deck_id_for_deterministic_output(tmp_path) -> None:
    pipeline = Pipeline(template_dir=TEMPLATE_DIR)

    output_path = pipeline.run(_requirement(), tmp_path, deck_id="deck-cli-test")

    assert output_path == tmp_path / "deck-cli-test.pptx"


def test_pipeline_renderer_type_hint_uses_renderer_protocol() -> None:
    hints = get_type_hints(Pipeline.__init__)

    assert hints["renderer"] == Renderer | None


def test_pipeline_accepts_renderer_protocol_implementation(tmp_path) -> None:
    class StubRenderer:
        def __init__(self) -> None:
            self.received_spec: SlideSpec | None = None
            self.received_output_dir: Path | None = None

        def render(self, spec: SlideSpec, output_dir: Path) -> Path:
            self.received_spec = spec
            self.received_output_dir = output_dir
            return output_dir / f"{spec.deck_id}.stub"

    renderer = StubRenderer()
    pipeline = Pipeline(template_dir=TEMPLATE_DIR, renderer=renderer)

    output_path = pipeline.run(_requirement(), tmp_path, deck_id="deck-stub")

    assert output_path == tmp_path / "deck-stub.stub"
    assert renderer.received_spec is not None
    assert renderer.received_spec.deck_id == "deck-stub"
    assert renderer.received_output_dir == tmp_path


def test_pipeline_passes_llm_provider_to_slide_builder(tmp_path) -> None:
    class StubRenderer:
        def __init__(self) -> None:
            self.received_spec: SlideSpec | None = None

        def render(self, spec: SlideSpec, output_dir: Path) -> Path:
            self.received_spec = spec
            return output_dir / f"{spec.deck_id}.stub"

    provider = FakeLlmProvider(
        LlmGenerationResult(
            slides=[
                GeneratedSlideContent(
                    slide_id="slide-001",
                    message="Fake provider message",
                    bullets=["fake bullet"],
                )
            ]
        )
    )
    renderer = StubRenderer()
    pipeline = Pipeline(
        template_dir=TEMPLATE_DIR,
        llm_provider=provider,
        renderer=renderer,
    )

    pipeline.run(_requirement(), tmp_path, deck_id="deck-fake-provider")

    assert renderer.received_spec is not None
    assert renderer.received_spec.slides[0].message == "Fake provider message"
    assert renderer.received_spec.slides[0].bullets == ["fake bullet"]


def test_pipeline_rejects_slide_builder_and_llm_provider_together() -> None:
    provider = FakeLlmProvider(
        LlmGenerationResult(
            slides=[
                GeneratedSlideContent(
                    slide_id="slide-001",
                    message="Fake provider message",
                )
            ]
        )
    )

    with pytest.raises(
        ValueError,
        match="slide_builder and llm_provider cannot be used together",
    ):
        Pipeline(
            template_dir=TEMPLATE_DIR,
            slide_builder=SlideBuilder(),
            llm_provider=provider,
        )
