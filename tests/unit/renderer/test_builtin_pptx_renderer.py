import pytest
from pptx import Presentation

from consultdeck.models.slide_spec import LayoutType, Slide, SlideSpec
from consultdeck.renderer.builtin_pptx_renderer import BuiltinPptxRenderer


def _slide_spec() -> SlideSpec:
    return SlideSpec(
        deck_id="deck-test",
        title="DX推進",
        template_id="proposal_standard",
        slides=[
            Slide(
                slide_id="slide-001",
                title="表紙",
                message="DX推進提案",
                bullets=[],
                layout_type=LayoutType.TITLE,
            ),
            Slide(
                slide_id="slide-002",
                title="課題",
                message="現状課題を整理します。",
                bullets=["情報が分散している", "意思決定が遅い"],
                layout_type=LayoutType.CONTENT,
                notes="課題を簡潔に説明する",
            ),
            Slide(
                slide_id="slide-003",
                title="比較",
                message="選択肢を比較します。",
                bullets=["現行案", "改善案"],
                layout_type=LayoutType.TWO_COLUMN,
            ),
            Slide(
                slide_id="slide-004",
                title="補足",
                message="補足",
                bullets=[],
                layout_type=LayoutType.BLANK,
            ),
        ],
    )


def test_renderer_creates_pptx_file(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()

    output_path = renderer.render(_slide_spec(), tmp_path)

    assert output_path == tmp_path / "deck-test.pptx"
    assert output_path.exists()
    assert output_path.suffix == ".pptx"


def test_renderer_creates_output_dir_when_missing(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()
    output_dir = tmp_path / "nested" / "output"

    output_path = renderer.render(_slide_spec(), output_dir)

    assert output_dir.exists()
    assert output_path.exists()


def test_renderer_creates_one_pptx_slide_per_slide_spec_slide(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()

    output_path = renderer.render(_slide_spec(), tmp_path)

    presentation = Presentation(output_path)
    assert len(presentation.slides) == len(_slide_spec().slides)


def test_renderer_handles_minimum_layouts(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()

    output_path = renderer.render(_slide_spec(), tmp_path)

    presentation = Presentation(output_path)
    text_by_slide = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in presentation.slides
    ]
    assert "表紙" in text_by_slide[0]
    assert "課題" in text_by_slide[1]
    assert "情報が分散している" in text_by_slide[1]
    assert "比較" in text_by_slide[2]
    assert "現行案" in text_by_slide[2]
    assert text_by_slide[3] == ""


def test_renderer_writes_slide_notes_to_speaker_notes(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()

    output_path = renderer.render(_slide_spec(), tmp_path)

    presentation = Presentation(output_path)
    assert presentation.slides[1].notes_slide.notes_text_frame.text == (
        "課題を簡潔に説明する"
    )


def test_renderer_handles_missing_notes(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()

    output_path = renderer.render(_slide_spec(), tmp_path)

    presentation = Presentation(output_path)
    assert presentation.slides[0].notes_slide.notes_text_frame.text == ""


def test_empty_slide_spec_is_rejected_by_model() -> None:
    with pytest.raises(ValueError):
        SlideSpec(
            deck_id="deck-empty",
            title="Empty",
            template_id="proposal_standard",
            slides=[],
        )


@pytest.mark.parametrize(
    "deck_id",
    [
        "../escape",
        "..",
        "nested/deck",
        "nested\\deck",
    ],
)
def test_invalid_deck_id_is_rejected(deck_id: str) -> None:
    with pytest.raises(ValueError):
        SlideSpec(
            deck_id=deck_id,
            title="DX推進",
            template_id="proposal_standard",
            slides=[
                Slide(
                    slide_id="slide-001",
                    title="課題",
                    message="課題",
                    bullets=[],
                    layout_type=LayoutType.CONTENT,
                )
            ],
        )


def test_renderer_cannot_write_outside_output_dir_with_deck_id(tmp_path) -> None:
    renderer = BuiltinPptxRenderer()
    output_dir = tmp_path / "output"

    output_path = renderer.render(_slide_spec(), output_dir)

    assert output_path.resolve().parent == output_dir.resolve()
    assert not (tmp_path / "deck-test.pptx").exists()
